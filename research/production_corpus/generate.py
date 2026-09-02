#!/usr/bin/env python
"""Generate the EN/VI enterprise production benchmark corpus.

Why this exists
---------------
`research/hardcases/` is a *mechanism* corpus: fourteen documents, each
isolating one way extraction breaks. It is the right shape for regression
locks and the wrong shape for deciding what to build next, because it says
nothing about how often a mechanism occurs or what it co-occurs with.

Production failure ranking (mission §6) needs a corpus with a realistic
*mixture*: mostly ordinary documents, a minority of hard ones, and — the part
the mechanism corpus structurally cannot provide — documents where several
difficulties land on the same page. Real enterprise documents are not one
mechanism at a time. A scanned Vietnamese purchase order with a seal across a
borderless table is a single document, not three.

Design
------
Documents are *composed* rather than hand-written. One builder takes a
`DocSpec` of independent options (table style, stamp, watermark, scan,
resolution, column count, ...) and renders it. Combining difficulties is then
a matter of setting two flags, which is what makes §30 combination cases cheap
enough to be worth having.

Everything is synthetic. Company names, people, tax codes and amounts are
invented and deliberately implausible as real registrations. No enterprise
document, public or private, is copied. That is a hard constraint: this corpus
is committed to git, and the private corpus it stands in for is not.

Determinism
-----------
Seeded, with fixed dates and no timestamps in content, so regenerating on
another machine produces the same bytes and the manifest's hashes stay
meaningful.

    python research/production_corpus/generate.py
"""
from __future__ import annotations

import argparse
import hashlib
import json
import random
from dataclasses import dataclass, field
from pathlib import Path

import pymupdf

REPO_ROOT = Path(__file__).resolve().parents[2]

A4 = (595.0, 842.0)
VN_FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
VN_FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

SEED = 20260902


# --------------------------------------------------------------------------
# Synthetic content. Invented entities; any resemblance to a real registration
# is unintended and the identifiers are deliberately not well-formed.
# --------------------------------------------------------------------------

VI_ORG = "CÔNG TY TNHH THƯƠNG MẠI MINH QUANG"
EN_ORG = "Brightwater Trading Limited"
VI_ADDR = "Số 12, Đường Lê Lợi, Phường Bến Nghé, Quận 1, TP. Hồ Chí Minh"
EN_ADDR = "44 Harbour Street, Unit 7, Wellington Central"
VI_PERSON = "Nguyễn Thị Hương"
EN_PERSON = "Alice Fairbanks"
TAX_ID = "0312998877"

VI_HEADER = ["CỘNG HÒA XÃ HỘI CHỦ NGHĨA VIỆT NAM", "Độc lập - Tự do - Hạnh phúc"]

# Per document type: (title_vi, title_en, body_vi, body_en, key/value pairs)
DOC_BODIES: dict[str, dict[str, list[str]]] = {
    "contract": {
        "title_vi": ["HỢP ĐỒNG MUA BÁN HÀNG HÓA"],
        "title_en": ["GOODS PURCHASE AGREEMENT"],
        "vi": [
            "Điều 1. Đối tượng của hợp đồng",
            "Bên A đồng ý bán và Bên B đồng ý mua hàng hóa theo danh mục đính kèm.",
            "Điều 2. Giá trị hợp đồng và phương thức thanh toán",
            "Tổng giá trị hợp đồng là 1.250.000.000 đồng, đã bao gồm thuế giá trị gia tăng.",
            "Điều 3. Thời hạn giao hàng",
            "Bên A giao hàng trong vòng 30 ngày kể từ ngày hợp đồng có hiệu lực.",
        ],
        "en": [
            "Article 1. Subject of the agreement",
            "The Seller agrees to sell and the Buyer agrees to purchase the goods listed in Annex A.",
            "Article 2. Contract value and payment terms",
            "The total contract value is USD 84,500.00, inclusive of all applicable taxes.",
            "Article 3. Delivery",
            "The Seller shall deliver within 30 days of the effective date of this agreement.",
        ],
    },
    "invoice": {
        "title_vi": ["HÓA ĐƠN GIÁ TRỊ GIA TĂNG"],
        "title_en": ["TAX INVOICE"],
        "vi": [
            "Ký hiệu: MQ/2026E    Số: 0004417",
            "Ngày 14 tháng 03 năm 2026",
            "Đơn vị bán hàng: " + VI_ORG,
            "Mã số thuế: " + TAX_ID,
        ],
        "en": [
            "Invoice number: INV-2026-004417",
            "Issue date: 14 March 2026",
            "Supplier: " + EN_ORG,
            "VAT registration: GB " + TAX_ID,
        ],
    },
    "financial_report": {
        "title_vi": ["BÁO CÁO TÀI CHÍNH QUÝ I NĂM 2026"],
        "title_en": ["QUARTERLY FINANCIAL REPORT — Q1 2026"],
        "vi": [
            "Đơn vị tính: triệu đồng",
            "Báo cáo này chưa được kiểm toán và chỉ dùng cho mục đích quản trị nội bộ.",
            "Doanh thu thuần trong kỳ tăng 12,4% so với cùng kỳ năm trước.",
        ],
        "en": [
            "Reporting currency: thousands of USD",
            "These figures are unaudited and prepared for internal management purposes.",
            "Net revenue for the period increased 12.4 per cent year on year.",
        ],
    },
    "company_profile": {
        "title_vi": ["HỒ SƠ NĂNG LỰC CÔNG TY"],
        "title_en": ["COMPANY CAPABILITY PROFILE"],
        "vi": [
            "Thành lập năm 2011, công ty hoạt động trong lĩnh vực phân phối thiết bị công nghiệp.",
            "Đội ngũ hiện tại gồm 148 nhân sự tại ba chi nhánh.",
        ],
        "en": [
            "Founded in 2011, the company distributes industrial equipment across three regions.",
            "The current headcount is 148 staff across three branch offices.",
        ],
    },
    "business_license": {
        "title_vi": ["GIẤY CHỨNG NHẬN ĐĂNG KÝ DOANH NGHIỆP"],
        "title_en": ["CERTIFICATE OF INCORPORATION"],
        "vi": [
            "Tên công ty: " + VI_ORG,
            "Mã số doanh nghiệp: " + TAX_ID,
            "Địa chỉ trụ sở chính: " + VI_ADDR,
            "Người đại diện theo pháp luật: " + VI_PERSON,
            "Vốn điều lệ: 8.000.000.000 đồng",
        ],
        "en": [
            "Company name: " + EN_ORG,
            "Company number: " + TAX_ID,
            "Registered office: " + EN_ADDR,
            "Director: " + EN_PERSON,
            "Share capital: USD 350,000",
        ],
    },
    "certificate": {
        "title_vi": ["CHỨNG NHẬN HỆ THỐNG QUẢN LÝ CHẤT LƯỢNG"],
        "title_en": ["QUALITY MANAGEMENT SYSTEM CERTIFICATE"],
        "vi": [
            "Chứng nhận rằng hệ thống quản lý của " + VI_ORG,
            "phù hợp với các yêu cầu của tiêu chuẩn ISO 9001:2015.",
            "Hiệu lực đến ngày 30 tháng 06 năm 2028.",
        ],
        "en": [
            "This is to certify that the management system of " + EN_ORG,
            "conforms to the requirements of ISO 9001:2015.",
            "Valid until 30 June 2028.",
        ],
    },
    "form": {
        "title_vi": ["PHIẾU ĐỀ NGHỊ THANH TOÁN"],
        "title_en": ["PAYMENT REQUISITION FORM"],
        "vi": [
            "Họ và tên người đề nghị: " + VI_PERSON,
            "Bộ phận: Phòng Kế toán",
            "Số tiền đề nghị: 24.750.000 đồng",
            "Lý do: Thanh toán chi phí vận chuyển tháng 02",
        ],
        "en": [
            "Requested by: " + EN_PERSON,
            "Department: Finance",
            "Amount requested: USD 1,875.00",
            "Purpose: February freight charges",
        ],
    },
    "policy": {
        "title_vi": ["QUY CHẾ QUẢN LÝ CHI TIÊU NỘI BỘ"],
        "title_en": ["INTERNAL EXPENDITURE POLICY"],
        "vi": [
            "Điều 1. Phạm vi áp dụng",
            "Quy chế này áp dụng cho toàn bộ nhân sự thuộc công ty và các chi nhánh.",
            "Điều 2. Thẩm quyền phê duyệt",
            "Khoản chi trên 50.000.000 đồng phải được Tổng Giám đốc phê duyệt bằng văn bản.",
        ],
        "en": [
            "Section 1. Scope",
            "This policy applies to all employees of the company and its branch offices.",
            "Section 2. Approval authority",
            "Expenditure above USD 5,000 requires written approval from the Managing Director.",
        ],
    },
    "technical_report": {
        "title_vi": ["BÁO CÁO KIỂM ĐỊNH KỸ THUẬT THIẾT BỊ"],
        "title_en": ["EQUIPMENT INSPECTION TECHNICAL REPORT"],
        "vi": [
            "Thiết bị: Máy nén khí trục vít, công suất 75 kW",
            "Kết quả đo độ rung tại vị trí ổ trục: 2,8 mm/s (giới hạn 4,5 mm/s)",
            "Kết luận: Thiết bị đạt yêu cầu vận hành an toàn.",
        ],
        "en": [
            "Equipment: Screw-type air compressor, rated 75 kW",
            "Measured bearing vibration velocity: 2.8 mm/s (limit 4.5 mm/s)",
            "Conclusion: the equipment meets safe operating requirements.",
        ],
    },
    "purchase_order": {
        "title_vi": ["ĐƠN ĐẶT HÀNG"],
        "title_en": ["PURCHASE ORDER"],
        "vi": [
            "Số đơn hàng: PO-2026-0318",
            "Nhà cung cấp: " + VI_ORG,
            "Điều kiện giao hàng: DDP kho Bên mua",
        ],
        "en": [
            "Order number: PO-2026-0318",
            "Supplier: " + EN_ORG,
            "Delivery terms: DDP buyer's warehouse",
        ],
    },
    "scanned_admin": {
        "title_vi": ["CÔNG VĂN V/V XÁC NHẬN CÔNG NỢ"],
        "title_en": ["OFFICIAL LETTER — CONFIRMATION OF OUTSTANDING BALANCE"],
        "vi": [
            "Kính gửi: Phòng Tài chính - Kế toán",
            "Căn cứ biên bản đối chiếu ngày 28 tháng 02 năm 2026,",
            "công nợ phải thu tại thời điểm khóa sổ là 412.600.000 đồng.",
        ],
        "en": [
            "To: Finance and Accounting Department",
            "Further to the reconciliation minutes dated 28 February 2026,",
            "the receivable balance at cut-off is USD 31,240.00.",
        ],
    },
}

TABLE_HEADERS_VI = ["STT", "Mô tả hàng hóa", "Đơn vị", "Số lượng", "Thành tiền"]
TABLE_HEADERS_EN = ["No.", "Item description", "Unit", "Qty", "Amount"]
TABLE_ROWS_VI = [
    ["1", "Sản phẩm A-100", "Cái", "12", "18.000.000"],
    ["2", "Bộ lọc khí Mã 22B", "Bộ", "4", "6.400.000"],
    ["3", "Dịch vụ lắp đặt", "Gói", "1", "3.500.000"],
]
TABLE_ROWS_EN = [
    ["1", "Product A-100", "pcs", "12", "18,000.00"],
    ["2", "Air filter type 22B", "set", "4", "6,400.00"],
    ["3", "Installation service", "job", "1", "3,500.00"],
]


# --------------------------------------------------------------------------
# Spec
# --------------------------------------------------------------------------


@dataclass
class DocSpec:
    """One document. Difficulty options are independent and composable."""

    document_id: str
    doc_type: str
    language: str                      # en | vi | mixed
    fmt: str = "pdf"                   # pdf | docx | xlsx | pptx | png
    pages: int = 1

    # difficulty options
    table: str | None = None           # ruled | borderless | merged | tiny
    stamp: str | None = None           # text | table | transparent | boundary
    watermark: bool = False
    tiny_text: bool = False
    low_contrast: bool = False
    scanned: bool = False
    scan_dpi: int = 150
    columns: int = 1
    rotate: int = 0
    broken_cmap: bool = False
    checkboxes: bool = False
    repeat_headfoot: bool = False

    hard_case_labels: list[str] = field(default_factory=list)
    notes: str = ""


def _lang_pick(spec: DocSpec) -> str:
    return "vi" if spec.language in ("vi", "mixed") else "en"


# --------------------------------------------------------------------------
# PDF rendering primitives
# --------------------------------------------------------------------------


def _text(page, point, s, size=11, bold=False, color=(0, 0, 0), rotate=0):
    """Draw text. `rotate` accepts any angle: PyMuPDF's own `rotate=` argument
    is limited to right angles, so anything else goes through a morph matrix
    about the insertion point (a 45-degree watermark needs this)."""
    font = VN_FONT_BOLD if bold else VN_FONT
    kw = dict(fontsize=size, fontfile=font, fontname="dvb" if bold else "dv", color=color)
    p = pymupdf.Point(*point)
    if rotate % 90 == 0:
        page.insert_text(p, s, rotate=rotate % 360, **kw)
    else:
        page.insert_text(p, s, morph=(p, pymupdf.Matrix(rotate)), **kw)


def _frag(line: str, limit: int = 34) -> str:
    """A short distinctive fragment of `line` to assert on.

    Asserting whole sentences makes a benchmark brittle for reasons that have
    nothing to do with extraction quality: a wrapped line, a hyphenation, or
    one clipped character fails the whole string. The existing hardcase corpus
    asserts short phrases for the same reason.
    """
    if len(line) <= limit:
        return line
    cut = line[:limit]
    return cut[:cut.rfind(" ")] if " " in cut else cut


def _draw_table(page, x, y, spec: DocSpec, lang: str) -> list[str]:
    """Draw a table in the style `spec.table` asks for. Returns strings that
    must survive extraction."""
    headers = TABLE_HEADERS_VI if lang == "vi" else TABLE_HEADERS_EN
    rows = TABLE_ROWS_VI if lang == "vi" else TABLE_ROWS_EN
    style = spec.table

    size = 4.2 if style == "tiny" else (7.5 if spec.tiny_text else 9.0)
    rh = 9.0 if style == "tiny" else 20.0
    widths = [26, 168, 44, 42, 92] if style != "tiny" else [16, 104, 28, 26, 58]
    ruled = style in ("ruled", "merged", "tiny")
    colour = (0.62, 0.62, 0.62) if spec.low_contrast else (0.15, 0.15, 0.15)

    total_w = sum(widths)
    n_rows = len(rows) + 1
    if ruled:
        for i in range(n_rows + 1):
            page.draw_line(pymupdf.Point(x, y + i * rh), pymupdf.Point(x + total_w, y + i * rh),
                           color=(0.3, 0.3, 0.3), width=0.5)
        cx = x
        for w in widths + [0]:
            page.draw_line(pymupdf.Point(cx, y), pymupdf.Point(cx, y + n_rows * rh),
                           color=(0.3, 0.3, 0.3), width=0.5)
            cx += w

    must: list[str] = []
    cx = x
    for w, h in zip(widths, headers):
        _text(page, (cx + 2, y + rh - 3), h, size=size, bold=True, color=colour)
        cx += w
    must.append(headers[1])

    for r, row in enumerate(rows, start=1):
        cx = x
        for c, (w, cell) in enumerate(zip(widths, row)):
            # A merged header spanning the two rightmost columns is the classic
            # native-grid failure: the finder keys on ruling lines and the
            # spanning cell has none beneath it.
            if style == "merged" and r == 1 and c == 3:
                _text(page, (cx + 2, y + r * rh + rh - 5),
                      "Số lượng / Thành tiền" if lang == "vi" else "Qty / Amount",
                      size=size, color=colour)
                must.append("Số lượng / Thành tiền" if lang == "vi" else "Qty / Amount")
                break
            _text(page, (cx + 2, y + r * rh + rh - 5), cell, size=size, color=colour)
            cx += w
        if style != "merged" or r != 1:
            must.append(row[1])
    return must


def _draw_stamp(page, centre, lang: str, transparent: bool = False):
    """A circular seal. Opaque fill removes the pixels underneath; the PDF text
    layer keeps them, which is the whole point of the occlusion cases."""
    red = (0.75, 0.05, 0.10)
    c = pymupdf.Point(*centre)
    fill = None if transparent else (0.95, 0.80, 0.82)
    page.draw_circle(c, 56, color=red, fill=fill, width=0)
    for r in (58, 52):
        page.draw_circle(c, r, color=red, width=2.4)
    _text(page, (c.x - 46, c.y - 4), "CÔNG TY TNHH" if lang == "vi" else "CERTIFIED",
          size=8, bold=True, color=red)
    _text(page, (c.x - 30, c.y + 12), "ĐÃ DUYỆT" if lang == "vi" else "APPROVED",
          size=8, bold=True, color=red)
    return ["ĐÃ DUYỆT" if lang == "vi" else "APPROVED"]


def _rasterize(doc, dpi: int, noisy: bool = True):
    """Flatten to images: the text layer disappears, forcing the visual route.
    This is how a scan differs from a born-digital PDF, and the reason the
    routing decision matters at all."""
    out = pymupdf.open()
    rng = random.Random(SEED)
    for page in doc:
        pix = page.get_pixmap(dpi=dpi)
        img = pix.tobytes("png")
        if noisy:
            # Mild JPEG-style degradation via re-encode at reduced quality.
            from io import BytesIO

            from PIL import Image
            im = Image.open(BytesIO(img)).convert("L")
            buf = BytesIO()
            im.save(buf, format="JPEG", quality=rng.choice([58, 62, 66]))
            img = buf.getvalue()
        p = out.new_page(width=page.rect.width, height=page.rect.height)
        p.insert_image(p.rect, stream=img)
    doc.close()
    return out


FIXED_DATE = "D:20260902000000Z"


def _save_pdf(doc, path: Path) -> None:
    """Save deterministically.

    PyMuPDF stamps a fresh /ID and a CreationDate on every save, so two runs of
    an otherwise identical generator produce different bytes. That silently
    makes the manifest's per-file sha256 useless — it would record which run
    happened, not which corpus. Pinning the metadata and suppressing the new
    /ID is what lets a regenerated corpus be checked against the one a result
    was measured on.
    """
    try:
        doc.subset_fonts()
    except Exception:  # noqa: BLE001 - optimization only
        pass
    doc.set_metadata({"producer": "doc-extraction production corpus generator",
                      "creator": "generate.py", "title": "", "author": "",
                      "subject": "", "keywords": "",
                      "creationDate": FIXED_DATE, "modDate": FIXED_DATE})
    doc.save(path, garbage=4, deflate=True, no_new_id=True)
    doc.close()


def build_pdf(spec: DocSpec, out_dir: Path) -> tuple[str, int, list[str], list[str]]:
    """Render one PDF. Returns (filename, page_count, must_contain, must_not_contain)."""
    if spec.broken_cmap:
        import sys
        sys.path.insert(0, str(REPO_ROOT))
        from tests.fixtures import make_pdf_with_broken_cmap
        path = out_dir / f"{spec.document_id}.pdf"
        make_pdf_with_broken_cmap(path, n_pages=spec.pages)
        return path.name, spec.pages, ["Cong Hoa", "Doanh Nghiep"], ["ЅoҖѸ", "Йoa"]

    lang = _lang_pick(spec)
    body = DOC_BODIES[spec.doc_type]
    title = body["title_vi"][0] if lang == "vi" else body["title_en"][0]
    lines = body[lang]

    doc = pymupdf.open()
    must: list[str] = []
    must_not: list[str] = []

    for pno in range(spec.pages):
        page = doc.new_page(width=A4[0], height=A4[1])
        y = 60.0

        if spec.watermark:
            _text(page, (110, 470), "BẢN SAO" if lang == "vi" else "COPY",
                  size=64, bold=True, color=(0.86, 0.86, 0.90), rotate=45)

        if lang == "vi" and spec.doc_type in ("business_license", "scanned_admin", "form"):
            for h in VI_HEADER:
                _text(page, (150, y), h, size=10, bold=True)
                y += 16
            must.append(VI_HEADER[1])
            y += 8

        if spec.repeat_headfoot:
            _text(page, (60, 36), VI_ORG if lang == "vi" else EN_ORG, size=8,
                  color=(0.4, 0.4, 0.4))
            _text(page, (60, A4[1] - 30), f"Trang {pno + 1}/{spec.pages}" if lang == "vi"
                  else f"Page {pno + 1} of {spec.pages}", size=8, color=(0.4, 0.4, 0.4))

        if pno == 0:
            _text(page, (60, y), title, size=14, bold=True)
            must.append(title)
            y += 30

        text_colour = (0.60, 0.60, 0.60) if spec.low_contrast else (0, 0, 0)
        body_size = 4.5 if spec.tiny_text else 11.0
        step = 11.0 if spec.tiny_text else 22.0

        if spec.columns == 2:
            # Columns are drawn with insert_textbox so long lines WRAP inside
            # the column. insert_text does not wrap, and text running past the
            # page edge is silently clipped out of the text layer — which
            # would make the ground truth unreachable rather than hard.
            half = (len(lines) + 1) // 2
            col_h = max(len(lines) * step, 160)
            for cx, chunk in ((60, lines[:half]), (315, lines[half:])):
                if not chunk:
                    continue
                page.insert_textbox(
                    pymupdf.Rect(cx, y, cx + 220, y + col_h), "\n".join(chunk),
                    fontsize=body_size, fontfile=VN_FONT, fontname="dv",
                    color=text_colour, align=0)
            must.append(_frag(lines[0]))
            if len(lines) > half:
                must.append(_frag(lines[half]))
            y += col_h + 12
        elif spec.rotate:
            # Rotated 90 degrees, text runs bottom-to-top from the insertion
            # point, so it must start low enough on the page to fit.
            for i, ln in enumerate(lines):
                _text(page, (60 + i * 24, A4[1] - 60), ln, size=body_size,
                      color=text_colour, rotate=spec.rotate)
            must.extend(_frag(x) for x in lines[:2])
            y += 40
        else:
            for i, ln in enumerate(lines):
                _text(page, (60, y + i * step), ln, size=body_size, color=text_colour)
            must.extend(_frag(x) for x in lines[:2])
            y += len(lines) * step + 12

        if spec.checkboxes:
            opts = (["Đồng ý", "Không đồng ý", "Cần bổ sung"] if lang == "vi"
                    else ["Approved", "Rejected", "More information required"])
            for i, o in enumerate(opts):
                bx = 62
                by = y + i * 20
                page.draw_rect(pymupdf.Rect(bx, by, bx + 10, by + 10),
                               color=(0.2, 0.2, 0.2), width=0.8)
                if i == 0:
                    page.draw_line(pymupdf.Point(bx + 1, by + 1),
                                   pymupdf.Point(bx + 9, by + 9), color=(0, 0, 0), width=1.2)
                    page.draw_line(pymupdf.Point(bx + 9, by + 1),
                                   pymupdf.Point(bx + 1, by + 9), color=(0, 0, 0), width=1.2)
                _text(page, (bx + 18, by + 9), o, size=10, color=text_colour)
            must.append(opts[0])
            y += len(opts) * 20 + 10

        if spec.table and pno == 0:
            must.extend(_draw_table(page, 60, y, spec, lang))
            y += 120

        if spec.stamp and pno == 0:
            centre = {"text": (185, 250), "table": (250, y - 60),
                      "transparent": (185, 250), "boundary": (300, y - 90)}[spec.stamp]
            must.extend(_draw_stamp(page, centre, lang,
                                    transparent=(spec.stamp == "transparent")))

    if spec.scanned:
        doc = _rasterize(doc, dpi=spec.scan_dpi)

    path = out_dir / f"{spec.document_id}.pdf"
    n = doc.page_count
    _save_pdf(doc, path)
    # de-duplicate while preserving order
    must = list(dict.fromkeys(must))
    return path.name, n, must, must_not


# --------------------------------------------------------------------------
# Office and image formats
# --------------------------------------------------------------------------


def _normalize_zip(path: Path) -> None:
    """Rewrite an OOXML file with fixed entry timestamps.

    Pinning core properties is not enough: docx/xlsx/pptx are zips, and each
    zip entry carries its own modification time from the moment it was
    written. Those bytes change every run, so the file hash does too. This
    rewrites the archive with a constant timestamp and constant entry order.
    """
    import re
    import shutil
    import zipfile
    stamp = (2026, 9, 2, 0, 0, 0)
    iso = "2026-09-02T00:00:00Z"
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
            tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for name in sorted(src.namelist()):
            info = zipfile.ZipInfo(name, date_time=stamp)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            data = src.read(name)
            if name == "docProps/core.xml":
                # openpyxl (and Word) rewrite dcterms:modified at save time,
                # overriding the value pinned on the properties object, so the
                # timestamp has to be normalized in the serialized XML too.
                data = re.sub(rb">[^<>]*?Z</dcterms:(created|modified)>",
                              lambda m: b">" + iso.encode() + b"</dcterms:"
                              + m.group(1) + b">", data)
            dst.writestr(info, data)
    shutil.move(tmp, path)


def _pin_office(props) -> None:
    """Same determinism problem as PDFs: OOXML core properties carry created /
    modified timestamps that change on every run."""
    from datetime import datetime
    stamp = datetime(2026, 9, 2, 0, 0, 0)
    for attr in ("created", "modified", "last_printed"):
        try:
            setattr(props, attr, stamp)
        except Exception:  # noqa: BLE001 - not every backend exposes all three
            pass
    for attr, val in (("author", "generate.py"), ("last_modified_by", "generate.py"),
                      ("creator", "generate.py"), ("revision", 1)):
        try:
            setattr(props, attr, val)
        except Exception:  # noqa: BLE001
            pass


def build_docx(spec: DocSpec, out_dir: Path):
    from docx import Document
    lang = _lang_pick(spec)
    body = DOC_BODIES[spec.doc_type]
    title = body["title_vi"][0] if lang == "vi" else body["title_en"][0]
    lines = body[lang]

    d = Document()
    d.add_heading(title, level=1)
    for ln in lines:
        d.add_paragraph(ln)
    must = [title, lines[0]]
    if spec.table:
        headers = TABLE_HEADERS_VI if lang == "vi" else TABLE_HEADERS_EN
        rows = TABLE_ROWS_VI if lang == "vi" else TABLE_ROWS_EN
        t = d.add_table(rows=len(rows) + 1, cols=len(headers))
        t.style = "Table Grid"
        for c, h in enumerate(headers):
            t.cell(0, c).text = h
        for r, row in enumerate(rows, start=1):
            for c, cell in enumerate(row):
                t.cell(r, c).text = cell
        must.extend([headers[1], rows[0][1]])
    _pin_office(d.core_properties)
    path = out_dir / f"{spec.document_id}.docx"
    d.save(path)
    _normalize_zip(path)
    return path.name, 1, list(dict.fromkeys(must)), []


def build_xlsx(spec: DocSpec, out_dir: Path):
    from openpyxl import Workbook
    lang = _lang_pick(spec)
    body = DOC_BODIES[spec.doc_type]
    title = body["title_vi"][0] if lang == "vi" else body["title_en"][0]
    headers = TABLE_HEADERS_VI if lang == "vi" else TABLE_HEADERS_EN
    rows = TABLE_ROWS_VI if lang == "vi" else TABLE_ROWS_EN

    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws["A1"] = title
    ws.append([])
    ws.append(headers)
    for row in rows:
        ws.append(row)
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = body[lang][0]
    _pin_office(wb.properties)
    path = out_dir / f"{spec.document_id}.xlsx"
    wb.save(path)
    _normalize_zip(path)
    # one page per sheet in this project's office route
    return path.name, 2, [title, headers[1], rows[0][1], body[lang][0]], []


def build_pptx(spec: DocSpec, out_dir: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    lang = _lang_pick(spec)
    body = DOC_BODIES[spec.doc_type]
    title = body["title_vi"][0] if lang == "vi" else body["title_en"][0]
    lines = body[lang]

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8), Inches(3))
    tf = box.text_frame
    tf.text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(16)
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Contact" if lang == "en" else "Liên hệ"
    b2 = s2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8), Inches(2))
    b2.text_frame.text = (EN_ADDR if lang == "en" else VI_ADDR)
    _pin_office(prs.core_properties)
    path = out_dir / f"{spec.document_id}.pptx"
    prs.save(path)
    _normalize_zip(path)
    return path.name, 2, [title, lines[0], (EN_ADDR if lang == "en" else VI_ADDR)], []


def build_png(spec: DocSpec, out_dir: Path):
    """Render the PDF form of the document and keep only the pixels."""
    tmp = DocSpec(**{**spec.__dict__, "fmt": "pdf", "document_id": spec.document_id + "__tmp"})
    name, _n, must, must_not = build_pdf(tmp, out_dir)
    src = out_dir / name
    doc = pymupdf.open(src)
    pix = doc[0].get_pixmap(dpi=spec.scan_dpi)
    path = out_dir / f"{spec.document_id}.png"
    pix.save(path)
    doc.close()
    src.unlink()
    return path.name, 1, must, must_not


BUILDERS = {"pdf": build_pdf, "docx": build_docx, "xlsx": build_xlsx,
            "pptx": build_pptx, "png": build_png}


# --------------------------------------------------------------------------
# The corpus
# --------------------------------------------------------------------------


def corpus_specs() -> list[DocSpec]:
    S: list[DocSpec] = []

    def add(**kw):
        S.append(DocSpec(**kw))

    # -- Ordinary documents. A production corpus is mostly these; if the
    # -- ranking in §6 is to mean anything, easy documents must be represented
    # -- in proportion rather than curated away.
    ordinary = [
        ("contract", "vi"), ("contract", "en"),
        ("invoice", "vi"), ("invoice", "en"),
        ("financial_report", "vi"), ("financial_report", "en"),
        ("company_profile", "vi"), ("company_profile", "en"),
        ("business_license", "vi"), ("business_license", "en"),
        ("certificate", "vi"), ("certificate", "en"),
        ("policy", "vi"), ("policy", "en"),
        ("technical_report", "vi"), ("technical_report", "en"),
        ("purchase_order", "vi"), ("purchase_order", "en"),
        ("form", "vi"), ("form", "en"),
    ]
    for dt, lg in ordinary:
        add(document_id=f"ord_{dt}_{lg}", doc_type=dt, language=lg,
            table="ruled" if dt in ("invoice", "purchase_order", "financial_report") else None,
            hard_case_labels=[], notes="ordinary born-digital document")

    # office formats
    add(document_id="ord_contract_docx_en", doc_type="contract", language="en", fmt="docx")
    add(document_id="ord_policy_docx_vi", doc_type="policy", language="vi", fmt="docx")
    add(document_id="ord_invoice_docx_vi", doc_type="invoice", language="vi", fmt="docx",
        table="ruled")
    add(document_id="ord_finreport_xlsx_en", doc_type="financial_report", language="en",
        fmt="xlsx", table="ruled")
    add(document_id="ord_purchase_xlsx_vi", doc_type="purchase_order", language="vi",
        fmt="xlsx", table="ruled")
    add(document_id="ord_profile_pptx_en", doc_type="company_profile", language="en", fmt="pptx")
    add(document_id="ord_profile_pptx_vi", doc_type="company_profile", language="vi", fmt="pptx")
    add(document_id="ord_certificate_png_en", doc_type="certificate", language="en",
        fmt="png", scan_dpi=200, hard_case_labels=["scan_quality"],
        notes="image input; no text layer by construction")
    add(document_id="ord_invoice_png_vi", doc_type="invoice", language="vi", fmt="png",
        table="ruled", scan_dpi=200, hard_case_labels=["scan_quality", "table_detection"])

    # -- Single-mechanism hard cases -------------------------------------
    add(document_id="hc_tiny_text_vi", doc_type="policy", language="vi", tiny_text=True,
        hard_case_labels=["tiny_text"], notes="4.5pt body; ~11px at 200 DPI, below the floor")
    add(document_id="hc_tiny_text_en", doc_type="policy", language="en", tiny_text=True,
        hard_case_labels=["tiny_text"])
    add(document_id="hc_low_contrast_vi", doc_type="contract", language="vi", low_contrast=True,
        hard_case_labels=["low_contrast"])
    add(document_id="hc_scan_vi", doc_type="scanned_admin", language="vi", scanned=True,
        hard_case_labels=["scan_quality"], notes="rasterized + JPEG degraded; no text layer")
    add(document_id="hc_scan_en", doc_type="scanned_admin", language="en", scanned=True,
        hard_case_labels=["scan_quality"])
    add(document_id="hc_stamp_text_vi", doc_type="business_license", language="vi", stamp="text",
        hard_case_labels=["stamp", "occlusion"],
        notes="opaque seal; native keeps the covered text, visual cannot")
    add(document_id="hc_stamp_table_vi", doc_type="invoice", language="vi", table="ruled",
        stamp="table", hard_case_labels=["stamp", "occlusion", "table_structure"],
        notes="the case no single strategy recovered in experiment 008")
    add(document_id="hc_transparent_seal_vi", doc_type="business_license", language="vi",
        stamp="transparent", hard_case_labels=["stamp", "occlusion"],
        notes="outline-only seal; pixels underneath survive, unlike the opaque case")
    add(document_id="hc_watermark_vi", doc_type="contract", language="vi", watermark=True,
        hard_case_labels=["watermark"])
    add(document_id="hc_multicolumn_en", doc_type="policy", language="en", columns=2,
        hard_case_labels=["multi_column", "reading_order"])
    add(document_id="hc_rotation_vi", doc_type="certificate", language="vi", rotate=90,
        hard_case_labels=["rotation", "reading_order"])
    add(document_id="hc_borderless_en", doc_type="invoice", language="en", table="borderless",
        hard_case_labels=["borderless_table", "table_detection"])
    add(document_id="hc_merged_en", doc_type="invoice", language="en", table="merged",
        hard_case_labels=["merged_cells", "table_structure"])
    add(document_id="hc_tiny_cells_vi", doc_type="invoice", language="vi", table="tiny",
        hard_case_labels=["tiny_text", "table_structure"])
    add(document_id="hc_encoding_vi", doc_type="business_license", language="vi",
        broken_cmap=True, pages=2, hard_case_labels=["encoding"],
        notes="regression lock: silent garbage is worse than visible failure")
    add(document_id="hc_crosspage_vi", doc_type="policy", language="vi", pages=4,
        repeat_headfoot=True, hard_case_labels=["cross_page"])
    add(document_id="hc_checkbox_en", doc_type="form", language="en", checkboxes=True,
        hard_case_labels=["form", "checkbox"])
    add(document_id="hc_checkbox_vi", doc_type="form", language="vi", checkboxes=True,
        hard_case_labels=["form", "checkbox", "key_value"])

    # -- Combination cases (§30). Real documents stack difficulties. -----
    add(document_id="cmb_stamp_table_vi", doc_type="purchase_order", language="vi",
        table="merged", stamp="table",
        hard_case_labels=["stamp", "occlusion", "merged_cells", "table_structure"])
    add(document_id="cmb_tiny_table_en", doc_type="invoice", language="en", table="tiny",
        tiny_text=True, hard_case_labels=["tiny_text", "table_structure"])
    add(document_id="cmb_scan_tiny_vi", doc_type="policy", language="vi", tiny_text=True,
        scanned=True, hard_case_labels=["scan_quality", "tiny_text"],
        notes="the compound case targeted high-DPI recovery exists for")
    add(document_id="cmb_multicol_table_en", doc_type="financial_report", language="en",
        columns=2, table="ruled",
        hard_case_labels=["multi_column", "reading_order", "table_detection"])
    add(document_id="cmb_watermark_tiny_vi", doc_type="contract", language="vi",
        watermark=True, tiny_text=True, hard_case_labels=["watermark", "tiny_text"])
    add(document_id="cmb_lowcontrast_stamp_vi", doc_type="business_license", language="vi",
        low_contrast=True, stamp="text",
        hard_case_labels=["low_contrast", "stamp", "occlusion"])
    add(document_id="cmb_scan_stamp_table_vi", doc_type="invoice", language="vi",
        table="ruled", stamp="table", scanned=True,
        hard_case_labels=["scan_quality", "stamp", "occlusion", "table_structure"],
        notes="triple compound; no text layer AND occlusion AND structure")
    add(document_id="cmb_scan_multicol_en", doc_type="policy", language="en", columns=2,
        scanned=True, hard_case_labels=["scan_quality", "multi_column", "reading_order"])
    add(document_id="cmb_borderless_lowcontrast_en", doc_type="invoice", language="en",
        table="borderless", low_contrast=True,
        hard_case_labels=["borderless_table", "low_contrast", "table_detection"])
    add(document_id="cmb_stamp_boundary_vi", doc_type="invoice", language="vi", table="ruled",
        stamp="boundary", hard_case_labels=["stamp", "occlusion", "table_structure"],
        notes="seal crosses a cell boundary rather than sitting inside one")

    # -- Long document (§29) ---------------------------------------------
    add(document_id="long_policy_vi_60p", doc_type="policy", language="vi", pages=60,
        repeat_headfoot=True, hard_case_labels=["cross_page", "long_document"],
        notes="memory, model reuse and partial-failure behaviour over 60 pages")

    return S


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    h.update(path.read_bytes())
    return h.hexdigest()


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=str(Path(__file__).resolve().parent / "corpus"))
    args = ap.parse_args(argv)

    out = Path(args.out)
    out.mkdir(parents=True, exist_ok=True)
    random.seed(SEED)

    specs = corpus_specs()
    entries = []
    for spec in specs:
        fname, pages, must, must_not = BUILDERS[spec.fmt](spec, out)
        path = out / fname
        entries.append({
            "document_id": spec.document_id,
            "filename": fname,
            "format": spec.fmt,
            "language": spec.language,
            "document_type": spec.doc_type,
            "page_count": pages,
            "source": "synthetic (research/production_corpus/generate.py)",
            "hard_case_labels": spec.hard_case_labels,
            "difficulty": "easy" if not spec.hard_case_labels else (
                "hard" if len(spec.hard_case_labels) < 3 else "extreme"),
            "must_contain": must,
            "must_not_contain": must_not,
            "expected_tables": 1 if spec.table else 0,
            "sha256": sha256(path),
            "bytes": path.stat().st_size,
            "notes": spec.notes,
        })

    manifest = {
        "corpus": "enterprise-production-en-vi",
        "version": 1,
        "generated_by": "research/production_corpus/generate.py",
        "seed": SEED,
        "synthetic": True,
        "contains_real_enterprise_data": False,
        "languages": ["en", "vi"],
        "documents": len(entries),
        "pages_total": sum(e["page_count"] for e in entries),
        "purpose": (
            "Primary EN/VI production benchmark. Unlike research/hardcases (one "
            "mechanism per document), this corpus has a realistic mixture and "
            "documents that combine difficulties, so failure frequency and "
            "co-occurrence can be measured rather than guessed."
        ),
        "scoring": (
            "recall over must_contain, NFC-normalized, identical to "
            "research/hardcases/run_benchmark.py; must_not_contain catches "
            "confident garbage"
        ),
        "documents_list": entries,
    }
    (out / "manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    print(f"{len(entries)} documents -> {out}")
    by_fmt: dict[str, int] = {}
    for e in entries:
        by_fmt[e["format"]] = by_fmt.get(e["format"], 0) + 1
    print("by format:", by_fmt)
    print("pages total:", manifest["pages_total"])
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
