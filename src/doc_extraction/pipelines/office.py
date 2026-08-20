"""Step C — native extraction for DOCX / XLSX / PPTX.

Parses structure directly from each format's own object model instead of
routing through rendering/OCR — this preserves headings, lists, and tables
far better than rasterizing and re-detecting them would, and is much
cheaper. This is the whole point of not immediately sending everything
through OCR/VLM.

Pagination semantics
--------------------
Rendered page boundaries are a *layout* property, not a document-structure
property, and the OOXML formats differ in whether they define one:

* **DOCX** — page breaks are computed by the renderer from fonts, margins
  and widow/orphan rules. python-docx cannot know them without a layout
  engine. So a DOCX yields **one logical `Page`** with
  ``is_rendered_page=False`` and ``page_number=None`` on every element.
  It does *not* claim to be page 1: fabricating pagination the format does
  not define would silently corrupt any downstream analysis keyed on page
  number. Getting real DOCX pagination requires rendering (LibreOffice or
  Word), which this project deliberately does not require — see
  docs/backends.md.
* **XLSX** — one `Page` per worksheet. A sheet is a real, addressable unit,
  so ``page_number`` is the 1-based sheet position, but
  ``is_rendered_page=False`` because a sheet is not a printed page either.
* **PPTX** — one `Page` per slide, ``is_rendered_page=True``: a slide *is*
  the rendered unit, and its dimensions are real.

Other known limitations (deliberate, documented):
- DOCX merged-cell span detection relies on python-docx repeating the same
  underlying `<w:tc>` element across spanned grid positions; unusual/nested
  merges may not be captured perfectly.
"""
from __future__ import annotations

from pathlib import Path

import docx
import openpyxl
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from pptx import Presentation

from doc_extraction.schemas.element import BBox, Element, ElementType
from doc_extraction.schemas.page import Page
from doc_extraction.schemas.table import Cell, Table
from doc_extraction.utils.logging import StageLogger, noop_stage

BACKEND_NAME = "native-office"

# Safety cap so a spreadsheet with an inflated `max_row`/`max_column` (a
# known openpyxl quirk when formatting extends past the real data) can't
# blow up runtime/memory. Documented, not silently truncated: a warning is
# recorded on the stage context when this triggers.
_XLSX_MAX_CELLS = 500_000


def _heading_level(style_name: str) -> int:
    for token in style_name.split():
        if token.isdigit():
            return int(token)
    return 1


def _convert_docx_table(docx_table: DocxTable, table_id: str) -> Table:
    n_rows = len(docx_table.rows)
    n_cols = len(docx_table.columns)
    row_cell_ids = [[id(c._tc) for c in row.cells] for row in docx_table.rows]

    grid_claimed: set[tuple[int, int]] = set()
    cells: list[Cell] = []
    for r, row in enumerate(docx_table.rows):
        row_cells = list(row.cells)
        c = 0
        seen_in_row: set[int] = set()
        while c < len(row_cells):
            if (r, c) in grid_claimed:
                c += 1
                continue
            tc_id = id(row_cells[c]._tc)
            if tc_id in seen_in_row:
                c += 1
                continue
            seen_in_row.add(tc_id)

            col_span = 1
            while c + col_span < len(row_cells) and id(row_cells[c + col_span]._tc) == tc_id:
                col_span += 1

            row_span = 1
            rr = r + 1
            while rr < n_rows and c < len(row_cell_ids[rr]) and row_cell_ids[rr][c] == tc_id:
                row_span += 1
                rr += 1

            for dr in range(row_span):
                for dc in range(col_span):
                    grid_claimed.add((r + dr, c + dc))

            cells.append(
                Cell(
                    row=r,
                    col=c,
                    row_span=row_span,
                    col_span=col_span,
                    text=row_cells[c].text.strip(),
                    is_header=(r == 0),
                )
            )
            c += col_span

    return Table(
        id=table_id,
        page_number=None,  # DOCX has no rendered pagination — see module docstring
        n_rows=n_rows,
        n_cols=n_cols,
        cells=cells,
        source_backend=BACKEND_NAME,
        confidence=1.0,
    )


def parse_docx(path: Path, logger: StageLogger | None = None) -> list[Page]:
    ctx_manager = logger.stage("parse", BACKEND_NAME) if logger else noop_stage()
    with ctx_manager as ctx:
        document = docx.Document(str(path))
        elements: list[Element] = []
        tables: list[Table] = []
        order_index = 0
        table_index = 0

        for child in document.element.body.iterchildren():
            if child.tag.endswith("}p"):
                para = Paragraph(child, document)
                text = para.text.strip()
                if not text:
                    continue
                style = (para.style.name or "").lower() if para.style is not None else ""
                if style.startswith("heading") or style.startswith("title"):
                    etype = ElementType.HEADING
                    level = _heading_level(style)
                elif "list" in style or "bullet" in style:
                    etype = ElementType.LIST_ITEM
                    level = None
                else:
                    etype = ElementType.PARAGRAPH
                    level = None
                elements.append(
                    Element(
                        id=f"p0-e{order_index}",
                        type=etype,
                        text=text,
                        page_number=None,  # no rendered pagination in DOCX
                        confidence=1.0,
                        source_backend=BACKEND_NAME,
                        order_index=order_index,
                        level=level,
                    )
                )
                order_index += 1
            elif child.tag.endswith("}tbl"):
                docx_table = DocxTable(child, document)
                table = _convert_docx_table(docx_table, f"p0-t{table_index}")
                tables.append(table)
                elements.append(
                    Element(
                        id=f"p0-e{order_index}",
                        type=ElementType.TABLE,
                        page_number=None,  # no rendered pagination in DOCX
                        confidence=1.0,
                        source_backend=BACKEND_NAME,
                        table_id=table.id,
                        order_index=order_index,
                    )
                )
                order_index += 1
                table_index += 1

        ctx.metrics = {"num_elements": len(elements), "num_tables": len(tables)}

    return [
        Page(
            index=0,
            width=0,
            height=0,
            is_rendered_page=False,  # a DOCX body is a logical unit, not a rendered page
            coordinate_unit="none",  # python-docx exposes no geometry for body content
            source_route="native_office",
            source_backend=BACKEND_NAME,
            elements=elements,
            tables=tables,
            notes=[
                "DOCX pagination is a renderer property and is not available without a "
                "layout engine; this document is represented as one logical page with "
                "page_number=None on every element."
            ],
        )
    ]


def _convert_sheet(sheet, table_id: str, page_number: int) -> tuple[Table, int, list[str]]:
    warnings: list[str] = []
    max_row = sheet.max_row or 0
    max_col = sheet.max_column or 0
    if max_row * max_col > _XLSX_MAX_CELLS:
        original = (max_row, max_col)
        # Keep the aspect ratio roughly intact while capping total cells.
        scale = (_XLSX_MAX_CELLS / (max_row * max_col)) ** 0.5
        max_row = max(1, int(max_row * scale))
        max_col = max(1, int(max_col * scale))
        warnings.append(
            f"sheet '{sheet.title}' used range {original[0]}x{original[1]} exceeds the "
            f"{_XLSX_MAX_CELLS}-cell safety cap; truncated to {max_row}x{max_col}"
        )

    merged_span: dict[tuple[int, int], tuple[int, int]] = {}
    covered: set[tuple[int, int]] = set()
    for merged_range in sheet.merged_cells.ranges:
        r0, c0 = merged_range.min_row - 1, merged_range.min_col - 1
        r1, c1 = merged_range.max_row - 1, merged_range.max_col - 1
        merged_span[(r0, c0)] = (r1 - r0 + 1, c1 - c0 + 1)
        for r in range(r0, r1 + 1):
            for c in range(c0, c1 + 1):
                if (r, c) != (r0, c0):
                    covered.add((r, c))

    cells: list[Cell] = []
    for r in range(max_row):
        for c in range(max_col):
            if (r, c) in covered:
                continue
            row_span, col_span = merged_span.get((r, c), (1, 1))
            value = sheet.cell(row=r + 1, column=c + 1).value
            if value is None and (row_span, col_span) == (1, 1):
                continue
            cells.append(
                Cell(
                    row=r,
                    col=c,
                    row_span=row_span,
                    col_span=col_span,
                    text="" if value is None else str(value),
                    is_header=(r == 0),
                )
            )

    table = Table(
        id=table_id,
        page_number=page_number,
        n_rows=max_row,
        n_cols=max_col,
        cells=cells,
        source_backend=BACKEND_NAME,
        confidence=1.0,
    )
    return table, len(cells), warnings


def parse_xlsx(path: Path, logger: StageLogger | None = None) -> list[Page]:
    ctx_manager = logger.stage("parse", BACKEND_NAME) if logger else noop_stage()
    with ctx_manager as ctx:
        workbook = openpyxl.load_workbook(str(path), data_only=True)
        pages: list[Page] = []
        total_cells = 0
        all_warnings: list[str] = []
        for sheet_index, sheet_name in enumerate(workbook.sheetnames):
            sheet = workbook[sheet_name]
            table, n_cells, warnings = _convert_sheet(
                sheet, table_id=f"p{sheet_index}-t0", page_number=sheet_index + 1
            )
            total_cells += n_cells
            all_warnings.extend(warnings)
            element = Element(
                id=f"p{sheet_index}-e0",
                type=ElementType.TABLE,
                page_number=sheet_index + 1,
                confidence=1.0,
                source_backend=BACKEND_NAME,
                table_id=table.id,
                order_index=0,
                extra={"sheet_name": sheet_name},
            )
            pages.append(
                Page(
                    index=sheet_index,
                    width=0,
                    height=0,
                    is_rendered_page=False,  # a worksheet is not a printed page
                    coordinate_unit="none",  # cell geometry is grid-based, not spatial
                    source_route="native_office",
                    source_backend=BACKEND_NAME,
                    elements=[element],
                    tables=[table],
                    notes=[f"XLSX worksheet {sheet_index + 1}: {sheet_name!r}"] + warnings,
                )
            )
        ctx.warnings = all_warnings
        ctx.metrics = {"num_sheets": len(pages), "num_cells": total_cells}

    return pages


def _shape_bbox(shape) -> BBox | None:
    try:
        if None in (shape.left, shape.top, shape.width, shape.height):
            return None
        return BBox(
            x0=shape.left, y0=shape.top, x1=shape.left + shape.width, y1=shape.top + shape.height
        )
    except Exception:
        return None


def _convert_pptx_table(pptx_table, table_id: str, page_number: int) -> Table:
    n_rows = len(pptx_table.rows)
    n_cols = len(pptx_table.columns)
    cells: list[Cell] = []
    for r in range(n_rows):
        for c in range(n_cols):
            cell = pptx_table.cell(r, c)
            if cell.is_spanned and not cell.is_merge_origin:
                continue
            row_span = cell.span_height if cell.is_merge_origin else 1
            col_span = cell.span_width if cell.is_merge_origin else 1
            cells.append(
                Cell(
                    row=r,
                    col=c,
                    row_span=row_span,
                    col_span=col_span,
                    text=(cell.text or "").strip(),
                    is_header=(r == 0),
                )
            )
    return Table(
        id=table_id,
        page_number=page_number,
        n_rows=n_rows,
        n_cols=n_cols,
        cells=cells,
        source_backend=BACKEND_NAME,
        confidence=1.0,
    )


def parse_pptx(path: Path, logger: StageLogger | None = None) -> list[Page]:
    ctx_manager = logger.stage("parse", BACKEND_NAME) if logger else noop_stage()
    with ctx_manager as ctx:
        presentation = Presentation(str(path))
        pages: list[Page] = []
        total_elements = 0
        for slide_index, slide in enumerate(presentation.slides):
            elements: list[Element] = []
            tables: list[Table] = []
            order_index = 0
            # python-pptx builds a fresh proxy object on each `.title`
            # access, so `shape is slide.shapes.title` is never True. Compare
            # the underlying XML element instead, which is stable.
            title_shape = slide.shapes.title
            title_element = title_shape._element if title_shape is not None else None
            for shape in slide.shapes:
                bbox = _shape_bbox(shape)
                if getattr(shape, "has_table", False):
                    table = _convert_pptx_table(
                        shape.table, f"p{slide_index}-t{order_index}", page_number=slide_index + 1
                    )
                    tables.append(table)
                    elements.append(
                        Element(
                            id=f"p{slide_index}-e{order_index}",
                            type=ElementType.TABLE,
                            bbox=bbox,
                            page_number=slide_index + 1,
                            confidence=1.0,
                            source_backend=BACKEND_NAME,
                            table_id=table.id,
                            order_index=order_index,
                        )
                    )
                    order_index += 1
                elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    is_title = title_element is not None and shape._element is title_element
                    etype = ElementType.HEADING if is_title else ElementType.TEXT
                    elements.append(
                        Element(
                            id=f"p{slide_index}-e{order_index}",
                            type=etype,
                            text=shape.text_frame.text.strip(),
                            bbox=bbox,
                            page_number=slide_index + 1,
                            confidence=1.0,
                            source_backend=BACKEND_NAME,
                            order_index=order_index,
                            level=1 if etype == ElementType.HEADING else None,
                        )
                    )
                    order_index += 1
                elif shape.shape_type is not None and getattr(shape, "image", None) is not None:
                    elements.append(
                        Element(
                            id=f"p{slide_index}-e{order_index}",
                            type=ElementType.IMAGE,
                            bbox=bbox,
                            page_number=slide_index + 1,
                            confidence=1.0,
                            source_backend=BACKEND_NAME,
                            order_index=order_index,
                        )
                    )
                    order_index += 1
            total_elements += len(elements)
            pages.append(
                Page(
                    index=slide_index,
                    width=float(presentation.slide_width or 0),
                    height=float(presentation.slide_height or 0),
                    coordinate_unit="emu",
                    coordinate_origin="top-left",
                    is_rendered_page=True,  # a slide is the rendered unit
                    source_route="native_office",
                    source_backend=BACKEND_NAME,
                    elements=elements,
                    tables=tables,
                )
            )
        ctx.metrics = {"num_slides": len(pages), "num_elements": total_elements}

    return pages
