"""Programmatically-built test fixtures.

The real sample documents at the repo root are immutable experimental input:
tests may read them but must never write, rename, or copy-modify them. Where
the real corpus does not cover a case (it contains no PPTX and no standalone
image, and only one corrupt PDF), we synthesize a minimal file into pytest's
`tmp_path` instead.

Everything here writes only under a caller-supplied temporary directory.
"""
from __future__ import annotations

from pathlib import Path


def make_pdf_with_text(path: Path, pages: list[str], page_size: tuple[float, float] = (595, 842)) -> Path:
    """A born-digital PDF with a real, correctly-encoded text layer."""
    import pymupdf

    doc = pymupdf.open()
    for text in pages:
        page = doc.new_page(width=page_size[0], height=page_size[1])
        page.insert_text((72, 100), text, fontsize=11)
    doc.save(path)
    doc.close()
    return path


def make_pdf_with_table(path: Path) -> Path:
    """A born-digital PDF containing one ruled 3x3 table.

    Drawn with real vector ruling lines and text, which is what PyMuPDF's
    native table finder keys on — so this exercises the digital-PDF table
    path rather than a mock.
    """
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page(width=595, height=842)

    left, top, col_w, row_h, n_rows, n_cols = 72, 100, 120, 30, 3, 3
    for r in range(n_rows + 1):
        y = top + r * row_h
        page.draw_line(pymupdf.Point(left, y), pymupdf.Point(left + n_cols * col_w, y))
    for c in range(n_cols + 1):
        x = left + c * col_w
        page.draw_line(pymupdf.Point(x, top), pymupdf.Point(x, top + n_rows * row_h))

    headers = ["Item", "Qty", "Price"]
    rows = [["Widget", "2", "10.00"], ["Gadget", "5", "25.50"]]
    for c, label in enumerate(headers):
        page.insert_text((left + c * col_w + 5, top + 20), label, fontsize=10)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            page.insert_text((left + c * col_w + 5, top + r * row_h + 20), value, fontsize=10)

    doc.save(path)
    doc.close()
    return path


def make_pdf_with_broken_cmap(path: Path, n_pages: int = 2) -> Path:
    """A PDF that reproduces the real corruption mode faithfully.

    The page draws ordinary Latin glyphs, but the font carries a deliberately
    wrong ``/ToUnicode`` CMap that maps a *subset* of letters onto Cyrillic
    codepoints. Extracted text is therefore plentiful and confidently wrong,
    with Latin and Cyrillic mixed inside single words — exactly the shape of
    the failure observed in the real corpus.

    Building it this way (rather than writing pre-garbled characters) matters:
    a synthetic PDF written with garbled *source* text silently loses those
    characters, because the base-14 fonts PyMuPDF can embed have no Cyrillic
    glyphs. Corrupting the mapping instead of the input reproduces the true
    mechanism and needs no font coverage at all.
    """
    import pymupdf

    # Map consonants to Cyrillic and leave vowels/digits alone, so decoded
    # words mix scripts instead of turning uniformly Cyrillic.
    letters = "BCDFGHJKLMNPQRSTVWXZbcdfghjklmnpqrstvwxz"
    pairs = [(ord(ch), 0x0400 + (i * 5) % 0x100) for i, ch in enumerate(letters)]
    entries = "\n".join(f"<{code:02X}> <{target:04X}>" for code, target in pairs)
    cmap = (
        "/CIDInit /ProcSet findresource begin\n12 dict begin\nbegincmap\n"
        "/CMapName /BrokenTest def\n/CMapType 2 def\n"
        "1 begincodespacerange\n<00> <FF>\nendcodespacerange\n"
        f"{len(pairs)} beginbfchar\n{entries}\nendbfchar\n"
        "endcmap\nCMapName currentdict /CMap defineresource pop\nend\nend"
    ).encode()

    lines = [
        "Cong Hoa Xa Hoi Chu Nghia Viet Nam Doc lap Tu do Hanh phuc",
        "Giay Chung Nhan Dang Ky Doanh Nghiep Cong Ty Trach Nhiem",
        "Ma so doanh nghiep 0000000000 Dang ky lan dau ngay thang nam",
        "Ten cong ty viet bang tieng Viet Dia chi tru so chinh Dien thoai",
        "Von dieu le mot tram trieu dong Danh sach thanh vien gop von",
        "Nguoi dai dien theo phap luat cua cong ty Giam doc Quoc tich",
    ]

    doc = pymupdf.open()
    for _ in range(n_pages):
        page = doc.new_page(width=595, height=842)
        for row, line in enumerate(lines):
            page.insert_text((72, 100 + row * 24), line, fontsize=11, fontname="helv")

    for page_index in range(n_pages):
        for font in doc.get_page_fonts(page_index):
            font_xref = font[0]
            cmap_xref = doc.get_new_xref()
            doc.update_object(cmap_xref, "<<>>")
            doc.update_stream(cmap_xref, cmap)
            doc.xref_set_key(font_xref, "ToUnicode", f"{cmap_xref} 0 R")

    doc.save(path)
    doc.close()
    return path


def make_empty_pdf(path: Path, n_pages: int = 1) -> Path:
    """A PDF with pages but no text layer at all — the 'scanned' shape,
    without needing a real scan."""
    import pymupdf

    doc = pymupdf.open()
    for _ in range(n_pages):
        doc.new_page(width=595, height=842)
    doc.save(path)
    doc.close()
    return path


def make_image(path: Path, size: tuple[int, int] = (400, 300), color: str = "white") -> Path:
    from PIL import Image

    Image.new("RGB", size, color).save(path)
    return path


def make_docx(path: Path) -> Path:
    import docx

    document = docx.Document()
    document.add_heading("Test Heading", level=1)
    document.add_paragraph("First body paragraph.")
    document.add_paragraph("Bullet item", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    document.save(path)
    return path


def make_xlsx(path: Path, n_sheets: int = 2) -> Path:
    import openpyxl

    workbook = openpyxl.Workbook()
    default = workbook.active
    default.title = "Sheet1"
    default["A1"], default["B1"] = "Name", "Value"
    default["A2"], default["B2"] = "alpha", 1
    for i in range(1, n_sheets):
        sheet = workbook.create_sheet(f"Sheet{i + 1}")
        sheet["A1"] = f"sheet {i + 1}"
    workbook.save(path)
    return path


def make_pptx(path: Path, n_slides: int = 2) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for i in range(n_slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Slide {i + 1} title"
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"Body text on slide {i + 1}"
    presentation.save(path)
    return path


def make_corrupt_pdf(path: Path) -> Path:
    """A file claiming to be a PDF that no parser can open."""
    path.write_bytes(b"%PDF-1.7\nthis is not actually a valid PDF body\n%%EOF\n")
    return path


# Text samples for text-quality tests. The "garbled" sample is copied from
# the real observed failure in FROGSLEAP_BUSINESS LICENSE.pdf (a broken
# ToUnicode CMap), so the tests exercise the actual failure mode rather than
# an invented one.
CLEAN_VIETNAMESE_TEXT = (
    "CỘNG HÒA XÃ HỘI CHỦ NGHĨA VIỆT NAM Độc lập Tự do Hạnh phúc "
    "GIẤY CHỨNG NHẬN ĐĂNG KÝ DOANH NGHIỆP CÔNG TY TRÁCH NHIỆM HỮU HẠN "
    "Mã số doanh nghiệp Đăng ký lần đầu ngày tháng năm Tên công ty viết bằng "
    "tiếng Việt Địa chỉ trụ sở chính Điện thoại Vốn điều lệ một trăm triệu đồng "
    "Danh sách thành viên góp vốn Người đại diện theo pháp luật của công ty"
)

GARBLED_CMAP_TEXT = (
    "&Ӝ1* +Ñ$ ;\xad +Ӝ, &+Ӫ 1*+Ƭ$ 9,ӊ7 1$0 ĈӝF OұS 7ӵ GR +ҥQK SK~F "
    "*,Ҩ< &+Ӭ1* 1+Ұ1 ĈĂ1* .é '2$1+ 1*+,ӊ3 &Ð1* 7< 75È&+ 1+,ӊ0 +Ӳ8 +Ҥ1 "
    "0m Vӕ GRDQK QJKLӋS ĈăQJ Nê O̯Q ÿ̯X QJj\\ WKiQJ QăP 7rQ F{QJ W\\ YLӃW EҵQJ "
    "WLӃQJ 9LӋW ĈӏD FKӍ WUө Vӣ FKtQK ĈLӋQ WKRҥL 9ӕQ ÿLӅX Oӊ Pӝt WUăP WULӋX ÿӗQJ"
)

CLEAN_ENGLISH_TEXT = (
    "This is an ordinary paragraph of English prose used to verify that the "
    "text quality assessor does not flag correctly decoded Latin text. It "
    "contains normal punctuation, a few numbers such as 2026 and 42, and "
    "sentences of a realistic length so that the word level statistics are "
    "meaningful rather than dominated by a handful of short tokens."
)
